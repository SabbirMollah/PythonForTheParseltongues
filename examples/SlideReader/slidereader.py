import pyttsx3
from pptx import Presentation

prs = Presentation("lecture.pptx")

text_to_read = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_to_read.append(run.text)

speaker = pyttsx3.init()
speaker.say(text_to_read)
speaker.setProperty('rate',100)
speaker.setProperty('volume', 0.9)
speaker.runAndWait()

#pip install pypiwin32